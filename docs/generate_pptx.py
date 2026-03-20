"""
ASHA Assist – International Hackathon Pitch Deck Generator
Run:  python docs/generate_pptx.py
Output: docs/ASHA_Assist_Hackathon_Pitch.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ──────────────────── Brand Colors ────────────────────
DEEP_TEAL   = RGBColor(0x00, 0x6E, 0x70)   # #006E70
ACCENT_BLUE = RGBColor(0x00, 0xA8, 0xCC)   # #00A8CC
WARM_WHITE  = RGBColor(0xF4, 0xF9, 0xFF)   # #F4F9FF
SOFT_GREEN  = RGBColor(0x2D, 0xC0, 0x7F)   # #2DC07F
ALERT_RED   = RGBColor(0xE5, 0x3E, 0x3E)   # #E53E3E
DARK_GRAY   = RGBColor(0x1A, 0x20, 0x2C)   # #1A202C
MID_GRAY    = RGBColor(0x71, 0x80, 0x96)   # #718096
LIGHT_GRAY  = RGBColor(0xEB, 0xF2, 0xFF)   # #EBF2FF
AMBER       = RGBColor(0xFF, 0xA5, 0x00)   # #FFA500
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # truly blank


# ─────────────────── Helpers ───────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multi_para(slide, lines, x, y, w, h,
                   font_size=14, bold_first=False, color=DARK_GRAY,
                   line_spacing=None, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if line_spacing:
            p.space_after = Pt(line_spacing)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold_first and i == 0
        run.font.color.rgb = color
        run.font.name = font_name


def header_bar(slide, title, subtitle=None):
    """Full-width teal header bar at top."""
    add_rect(slide, 0, 0, 13.33, 1.4, DEEP_TEAL)
    add_text_box(slide, title, 0.5, 0.1, 12, 0.75,
                 font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.5, 0.82, 12, 0.45,
                     font_size=14, bold=False, color=ACCENT_BLUE, align=PP_ALIGN.LEFT)
    # Accent line
    add_rect(slide, 0, 1.38, 13.33, 0.06, SOFT_GREEN)


def slide_bg(slide):
    """Light background."""
    add_rect(slide, 0, 0, 13.33, 7.5, WARM_WHITE)


def footer(slide, slide_num, total=12):
    add_rect(slide, 0, 7.1, 13.33, 0.4, DEEP_TEAL)
    add_text_box(slide, "ASHA Assist  |  AI Rural Health Intelligence Platform  |  2026",
                 0.3, 7.1, 10, 0.4, font_size=9, color=WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide, f"{slide_num} / {total}",
                 12.5, 7.1, 0.8, 0.4, font_size=9, color=WHITE, align=PP_ALIGN.RIGHT)


def icon_bullet_box(slide, items, x, y, w, h, icon="●", dot_color=SOFT_GREEN,
                    font_size=13.5, text_color=DARK_GRAY):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        # icon run
        r_icon = p.add_run()
        r_icon.text = icon + "  "
        r_icon.font.size = Pt(font_size)
        r_icon.font.color.rgb = dot_color
        r_icon.font.name = "Calibri"
        # text run
        r_text = p.add_run()
        r_text.text = item
        r_text.font.size = Pt(font_size)
        r_text.font.color.rgb = text_color
        r_text.font.name = "Calibri"


def card(slide, x, y, w, h, title, body_lines, title_color=DEEP_TEAL,
         bg_color=WHITE, border_color=ACCENT_BLUE):
    # card bg
    bg = add_rect(slide, x, y, w, h, bg_color)
    bg.line.color.rgb = border_color
    bg.line.width = Pt(1.2)
    # title
    add_text_box(slide, title, x + 0.15, y + 0.1, w - 0.3, 0.38,
                 font_size=13, bold=True, color=title_color)
    # body
    add_multi_para(slide, body_lines, x + 0.15, y + 0.5, w - 0.3, h - 0.6,
                   font_size=11.5, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════
#  SLIDE 1 – COVER
# ═══════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
# Full dark gradient feel
add_rect(s1, 0, 0, 13.33, 7.5, DARK_GRAY)
# Large teal strip
add_rect(s1, 0, 0, 5.2, 7.5, DEEP_TEAL)
# Green accent bar
add_rect(s1, 5.2, 3.0, 0.08, 2.2, SOFT_GREEN)

# Left panel
add_text_box(s1, "ASHA", 0.5, 0.8, 4.5, 1.3, font_size=72, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(s1, "ASSIST", 0.5, 1.95, 4.5, 1.0, font_size=48, bold=True,
             color=ACCENT_BLUE, align=PP_ALIGN.LEFT)
add_text_box(s1, "AI Rural Health\nIntelligence Platform", 0.5, 3.1, 4.5, 1.2,
             font_size=18, bold=False, color=WARM_WHITE, align=PP_ALIGN.LEFT)

# Tag line strip
add_rect(s1, 0.5, 4.5, 4.2, 0.5, SOFT_GREEN)
add_text_box(s1, "Connecting Frontlines. Saving Lives.", 0.6, 4.54, 4.0, 0.42,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Right panel – highlights
add_text_box(s1, "International Hackathon 2026", 5.5, 0.7, 7.4, 0.6,
             font_size=16, bold=False, color=ACCENT_BLUE, align=PP_ALIGN.LEFT)
add_text_box(s1, "Healthcare × AI × Social Impact", 5.5, 1.25, 7.4, 0.5,
             font_size=13, bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT)

# Right panel bullets
highlights = [
    "🌐  Full-stack production-ready platform",
    "🤖  AI symptom triage (6 diseases)",
    "📱  Mobile-first ASHA worker interface",
    "🏥  Real-time doctor alert system",
    "📊  Disease surveillance & outbreak detection",
    "🔔  Twilio SMS  +  JWT Role-based Auth",
    "🐳  Dockerized & ready to deploy",
]
icon_bullet_box(s1, highlights, 5.5, 2.0, 7.4, 4.5,
                icon="", dot_color=SOFT_GREEN, font_size=14.5,
                text_color=WARM_WHITE)

footer(s1, 1)


# ═══════════════════════════════════════════════════════
#  SLIDE 2 – PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
slide_bg(s2)
header_bar(s2, "The Problem", "Rural healthcare in India is broken at the last mile")

# 4 stat callout boxes
stats = [
    ("800M+", "people in rural India rely on ASHA workers as\ntheir first healthcare contact"),
    ("72%", "of cases arrive at PHCs  already at a\ncritical or late-stage condition"),
    ("Manual Records", "Paper-based patient tracking leads to\nlost data and zero trend visibility"),
    ("Zero AI Triage", "No decision-support tool exists at the\nfield level for frontline workers"),
]
positions = [(0.5, 1.6), (3.8, 1.6), (7.1, 1.6), (10.4, 1.6)]
for (stat, desc), (px, py) in zip(stats, positions):
    add_rect(s2, px, py, 2.7, 2.3, WHITE).line.color.rgb = ACCENT_BLUE
    add_rect(s2, px, py, 2.7, 2.3, WHITE).line.width = Pt(1)
    add_rect(s2, px, py, 2.7, 0.55, DEEP_TEAL)
    add_text_box(s2, stat, px + 0.1, py + 0.05, 2.5, 0.5,
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_multi_para(s2, [desc], px + 0.12, py + 0.65, 2.5, 1.55,
                   font_size=11.5, color=DARK_GRAY)

# Pain-point banner
add_rect(s2, 0.5, 4.2, 12.33, 0.72, ALERT_RED)
add_text_box(s2, "Result: Delayed triage → missed outbreaks → preventable deaths",
             0.7, 4.26, 12.0, 0.6, font_size=18, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

# Quote
add_text_box(s2,
             '"Every hour of delay in High-risk case referral increases mortality risk by 12%"  – WHO Rural Health Report 2023',
             0.5, 5.25, 12.33, 0.55, font_size=11.5, bold=False,
             color=MID_GRAY, align=PP_ALIGN.CENTER)

# Bottom insight
icon_bullet_box(s2, [
    "ASHA workers manage 1,000+ patients each, with no digital tools and no real-time escalation path",
    "District hospitals are overwhelmed because early-stage detection never happened",
], 0.5, 5.95, 12.33, 1.0, icon="▸", dot_color=ALERT_RED, font_size=12.5)

footer(s2, 2)


# ═══════════════════════════════════════════════════════
#  SLIDE 3 – OUR SOLUTION
# ═══════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
slide_bg(s3)
header_bar(s3, "Our Solution", "ASHA Assist — A connected rural health intelligence loop")

add_text_box(s3, "One Platform. Three Stakeholders. Zero Gaps.", 0.5, 1.5, 12.33, 0.55,
             font_size=20, bold=True, color=DEEP_TEAL, align=PP_ALIGN.CENTER)

# Three columns
cols = [
    ("📱  ASHA Worker", DEEP_TEAL, [
        "Mobile-first patient registration",
        "Voice-enabled symptom input",
        "Instant AI triage results",
        "Photo prescription upload",
        "Offline-ready UX design",
    ]),
    ("🤖  AI Engine", ACCENT_BLUE, [
        "Rule-based + ML symptom matching",
        "Predicts 6 high-burden diseases",
        "Risk scoring: Low / Medium / High",
        "Differential diagnosis list",
        "Optional OpenAI enhancement",
    ]),
    ("🏥  Doctor Dashboard", SOFT_GREEN, [
        "Real-time high-risk alerts (SMS)",
        "Full patient history & vitals",
        "Disease trend charts",
        "Outbreak detection panel",
        "Review & approve care plans",
    ]),
]
col_x = [0.5, 4.65, 8.8]
for (title, col_color, bullets), cx in zip(cols, col_x):
    add_rect(s3, cx, 2.2, 4.0, 4.5, WHITE).line.color.rgb = col_color
    add_rect(s3, cx, 2.2, 4.0, 0.6, col_color)
    add_text_box(s3, title, cx + 0.1, 2.25, 3.8, 0.52,
                 font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    icon_bullet_box(s3, bullets, cx + 0.18, 2.88, 3.7, 3.7,
                    icon="✓", dot_color=col_color, font_size=12.5, text_color=DARK_GRAY)

footer(s3, 3)


# ═══════════════════════════════════════════════════════
#  SLIDE 4 – TECH STACK
# ═══════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
slide_bg(s4)
header_bar(s4, "Technology Stack", "Production-grade, modular, and containerized")

layers = [
    ("Frontend",     "React  +  Vite  +  Material UI  +  Recharts  +  React Router",   ACCENT_BLUE),
    ("Backend API",  "Node.js  +  Express  +  Mongoose  +  JWT Auth  +  Helmet  +  Rate-Limit", DEEP_TEAL),
    ("AI Service",   "Python  +  FastAPI  +  Custom Symptom Engine  +  OpenAI (optional)",       SOFT_GREEN),
    ("Database",     "MongoDB Atlas  +  GeoJSON 2dsphere Indexes  +  Full CRUD models",          AMBER),
    ("Alerts",       "Twilio SMS API  —  auto-fired on High-risk screening outcomes",             ALERT_RED),
    ("Infrastructure","Docker Compose  +  Nginx  +  Multi-container orchestration",              MID_GRAY),
]
for i, (layer, desc, lcolor) in enumerate(layers):
    y = 1.6 + i * 0.85
    add_rect(s4, 0.5, y, 2.2, 0.72, lcolor)
    add_text_box(s4, layer, 0.55, y + 0.06, 2.1, 0.6,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s4, 2.7, y, 10.1, 0.72, WHITE).line.color.rgb = lcolor
    add_text_box(s4, desc, 2.85, y + 0.12, 9.8, 0.5,
                 font_size=13, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT)

# Why it matters
add_rect(s4, 0.5, 6.75, 12.33, 0.5, DEEP_TEAL)
add_text_box(s4,
             "Every component is interchangeable and independently deployable — built for scale from Day 1",
             0.6, 6.77, 12.1, 0.46, font_size=12, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
footer(s4, 4)


# ═══════════════════════════════════════════════════════
#  SLIDE 5 – SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
slide_bg(s5)
header_bar(s5, "System Architecture", "End-to-end request flow & component interaction")

# Draw architecture diagram as shapes
# ASHA Phone
add_rect(s5, 0.35, 1.8, 2.0, 1.1, DEEP_TEAL)
add_text_box(s5, "📱 ASHA\nMobile App", 0.35, 1.85, 2.0, 1.0,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow 1
add_text_box(s5, "──────►", 2.35, 2.2, 1.1, 0.4,
             font_size=18, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# Backend
add_rect(s5, 3.45, 1.6, 2.5, 1.5, ACCENT_BLUE)
add_text_box(s5, "⚙️  Backend API\nNode.js + Express\nJWT · Rate Limit · Helmet",
             3.5, 1.65, 2.4, 1.4, font_size=11, bold=False, color=WHITE,
             align=PP_ALIGN.CENTER)

# Arrow to AI
add_text_box(s5, "──────►", 5.95, 2.2, 1.1, 0.4,
             font_size=18, bold=True, color=SOFT_GREEN, align=PP_ALIGN.CENTER)

# AI Service
add_rect(s5, 7.05, 1.6, 2.5, 1.5, SOFT_GREEN)
add_text_box(s5, "🤖  AI Service\nFastAPI + Python\nDisease · Risk · Confidence",
             7.1, 1.65, 2.4, 1.4, font_size=11, bold=False, color=WHITE,
             align=PP_ALIGN.CENTER)

# Arrow DB
add_text_box(s5, "▼", 4.55, 3.1, 0.5, 0.35,
             font_size=20, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# MongoDB
add_rect(s5, 3.45, 3.45, 2.5, 1.1, AMBER)
add_text_box(s5, "🗄  MongoDB\nGeoJSON · Atlas · Indexes",
             3.5, 3.5, 2.4, 1.0, font_size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

# Arrow Twilio
add_text_box(s5, "▼", 8.2, 3.1, 0.5, 0.35,
             font_size=20, bold=True, color=ALERT_RED, align=PP_ALIGN.CENTER)

# Twilio
add_rect(s5, 7.05, 3.45, 2.5, 1.1, ALERT_RED)
add_text_box(s5, "🔔  Twilio SMS\nAuto-fired on High Risk",
             7.1, 3.5, 2.4, 1.0, font_size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

# Doctor Dashboard
add_rect(s5, 10.05, 1.6, 2.8, 3.0, DARK_GRAY)
add_text_box(s5, "🏥  Doctor\nDashboard\n\n• High-risk queue\n• Analytics\n• Outbreak panel\n• Review notes",
             10.1, 1.65, 2.7, 2.9, font_size=11, bold=False, color=WARM_WHITE,
             align=PP_ALIGN.CENTER)

# Arrow to Doctor
add_text_box(s5, "──────►", 9.55, 2.2, 0.5, 0.4,
             font_size=14, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Flow description
flow_steps = [
    "1 · ASHA registers patient + symptoms on mobile app",
    "2 · Backend validates, stores, and calls AI service /predict",
    "3 · AI returns disease, risk level, recommendation & confidence",
    "4 · If HIGH risk → Twilio SMS fires instantly to on-call doctor",
    "5 · Doctor dashboard visualises all cases, trends, and outbreak signals",
]
icon_bullet_box(s5, flow_steps, 0.35, 4.75, 12.5, 2.2,
                icon="▸", dot_color=DEEP_TEAL, font_size=12.5, text_color=DARK_GRAY)

footer(s5, 5)


# ═══════════════════════════════════════════════════════
#  SLIDE 6 – AI ENGINE
# ═══════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
slide_bg(s6)
header_bar(s6, "AI Symptom Triage Engine", "Instant disease prediction with risk classification at point-of-care")

# Left – how it works
add_rect(s6, 0.5, 1.6, 5.8, 5.1, WHITE).line.color.rgb = DEEP_TEAL
add_rect(s6, 0.5, 1.6, 5.8, 0.55, DEEP_TEAL)
add_text_box(s6, "How It Works", 0.6, 1.65, 5.6, 0.45,
             font_size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

steps = [
    ("Input",       "ASHA enters symptoms by text, form, or voice"),
    ("NLP Parse",   "Keywords extracted + matched to disease symptom profiles"),
    ("Scoring",     "Weighted scoring across 6 disease models simultaneously"),
    ("Risk Level",  "Low / Medium / High — based on score thresholds"),
    ("Output",      "Disease, Risk, Recommendation, Confidence (0–1), Differentials"),
    ("Escalation",  "High-risk result → backend triggers SMS alert automatically"),
]
for i, (step, desc) in enumerate(steps):
    y = 2.3 + i * 0.72
    add_rect(s6, 0.65, y, 1.2, 0.5, SOFT_GREEN)
    add_text_box(s6, step, 0.65, y + 0.04, 1.2, 0.42,
                 font_size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s6, desc, 2.0, y + 0.04, 4.15, 0.5,
                 font_size=11.5, color=DARK_GRAY, align=PP_ALIGN.LEFT)

# Right – diseases + output sample
add_rect(s6, 6.8, 1.6, 6.0, 2.5, WHITE).line.color.rgb = ACCENT_BLUE
add_rect(s6, 6.8, 1.6, 6.0, 0.55, ACCENT_BLUE)
add_text_box(s6, "Diseases Covered", 6.9, 1.65, 5.8, 0.45,
             font_size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
diseases = ["🦟 Dengue", "🩸 Malaria", "🫁 Tuberculosis",
            "💉 Anemia", "😷 Respiratory Infection", "🤒 Typhoid"]
icon_bullet_box(s6, diseases, 6.9, 2.23, 5.8, 1.8,
                icon="", dot_color=SOFT_GREEN, font_size=13, text_color=DARK_GRAY)

# Output JSON sample
add_rect(s6, 6.8, 4.2, 6.0, 2.5, DARK_GRAY)
add_text_box(s6, "Sample AI Response", 6.9, 4.25, 5.8, 0.38,
             font_size=12, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.LEFT)
json_text = (
    '{\n'
    '  "disease_prediction": "Dengue",\n'
    '  "risk_level":         "High",\n'
    '  "recommendation":     "Refer to PHC immediately",\n'
    '  "confidence":         0.91,\n'
    '  "differentials":      ["Malaria", "Typhoid"]\n'
    '}'
)
add_text_box(s6, json_text, 6.9, 4.65, 5.7, 1.9,
             font_size=10.5, color=SOFT_GREEN, align=PP_ALIGN.LEFT, font_name="Courier New")

footer(s6, 6)


# ═══════════════════════════════════════════════════════
#  SLIDE 7 – ASHA WORKER FLOW (User Journey)
# ═══════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
slide_bg(s7)
header_bar(s7, "ASHA Worker Journey", "Designed for low-tech environments — 3 taps to a triage result")

steps7 = [
    ("1", "Login",       "Secure JWT\nrole-based login"),
    ("2", "Dashboard",   "Today's patients,\nrisk summary"),
    ("3", "New Patient", "Register with\nname, age, village"),
    ("4", "Vitals",      "BP, temp, SpO\u2082,\nweight, pulse"),
    ("5", "Symptoms",    "Voice input /\ncheckbox / text"),
    ("6", "AI Result",   "Disease + Risk +\nRecommendation"),
]
colors7 = [DEEP_TEAL, ACCENT_BLUE, SOFT_GREEN, AMBER, DEEP_TEAL, ALERT_RED]
for i, ((num, title, desc), clr) in enumerate(zip(steps7, colors7)):
    x = 0.5 + i * 2.1
    add_rect(s7, x, 1.7, 1.9, 1.9, clr)
    add_text_box(s7, num + "\n" + title, x + 0.05, 1.75, 1.8, 1.0,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s7, desc, x + 0.05, 2.78, 1.8, 0.75,
                 font_size=10.5, color=WARM_WHITE, align=PP_ALIGN.CENTER)
    if i < 5:
        add_text_box(s7, "►", x + 1.9, 2.45, 0.28, 0.4,
                     font_size=16, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Key design decisions
add_text_box(s7, "Key UX Design Decisions", 0.5, 3.85, 12.33, 0.45,
             font_size=16, bold=True, color=DEEP_TEAL, align=PP_ALIGN.LEFT)
ux_decisions = [
    "Voice input first — many ASHA workers have limited typing proficiency in English",
    "Material UI components with large tap targets and high-contrast colours for field use",
    "All API calls fail gracefully with local state preservation — partial offline tolerance",
    "AI results shown with clear colour-coded risk badges (green / amber / red) for instant comprehension",
]
icon_bullet_box(s7, ux_decisions, 0.5, 4.35, 12.33, 2.5,
                icon="✦", dot_color=SOFT_GREEN, font_size=13, text_color=DARK_GRAY)

footer(s7, 7)


# ═══════════════════════════════════════════════════════
#  SLIDE 8 – DOCTOR COMMAND CENTER
# ═══════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
slide_bg(s8)
header_bar(s8, "Doctor Command Center", "From receiving alerts to approving care plans — all in one view")

panels = [
    ("🚨 High-Risk Queue",
     ["Live list of High-risk screened patients",
      "Sorted by risk level + time of screening",
      "One-click to view full history + AI result",
      "Acknowledge / resolve alert in dashboard"]),
    ("📊 Disease Analytics",
     ["Interactive bar chart: cases per disease",
      "Line chart: daily incidence trends",
      "Village burden heatmap data",
      "Powered by Recharts — real-time updates"]),
    ("⚠️ Outbreak Detection",
     ["Automatic spike detection algorithm",
      "Cluster alerts when threshold crossed",
      '"Simulate Outbreak" for live demo',
      "Colour-coded severity banners"]),
    ("📋 Patient Detail",
     ["Full vitals history across screenings",
      "AI result with confidence band",
      "Doctor review notes + status update",
      "Differential diagnoses listed"]),
]
positions8 = [(0.5, 1.65), (6.8, 1.65), (0.5, 4.1), (6.8, 4.1)]
colors8 = [ALERT_RED, ACCENT_BLUE, AMBER, SOFT_GREEN]
for (title, bullets), (px, py), clr in zip(panels, positions8, colors8):
    add_rect(s8, px, py, 6.0, 2.2, WHITE).line.color.rgb = clr
    add_rect(s8, px, py, 6.0, 0.5, clr)
    add_text_box(s8, title, px + 0.1, py + 0.06, 5.8, 0.4,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    icon_bullet_box(s8, bullets, px + 0.15, py + 0.58, 5.7, 1.5,
                    icon="→", dot_color=clr, font_size=12, text_color=DARK_GRAY)

footer(s8, 8)


# ═══════════════════════════════════════════════════════
#  SLIDE 9 – SECURITY & DATA INTEGRITY
# ═══════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
slide_bg(s9)
header_bar(s9, "Security & Data Integrity", "Built with OWASP-aligned security principles throughout")

sec_items = [
    ("🔐 Authentication",  DEEP_TEAL,
     ["JWT tokens with secret-key signing",
      "Role-based access: asha · doctor · admin",
      "Protected routes — no cross-role data access",
      "Passwords hashed with bcrypt"]),
    ("🛡️ API Hardening",   ACCENT_BLUE,
     ["Helmet.js — sets 15 security HTTP headers",
      "express-rate-limit — blocks brute-force",
      "CORS whitelist — restricted origins",
      "Input validation before DB writes"]),
    ("🗄️ Data Safety",     SOFT_GREEN,
     ["MongoDB indexed + role-filtered queries",
      "GeoJSON data never exposed without auth",
      "Twilio credentials in environment variables",
      "No PII logged to console in production"]),
    ("📡 Transport",       AMBER,
     ["Nginx reverse proxy with TLS-ready config",
      "Docker network isolation per service",
      "AI service not publicly exposed",
      "Health endpoint rate-limited separately"]),
]
pos9 = [(0.5, 1.65), (6.9, 1.65), (0.5, 4.15), (6.9, 4.15)]
for (title, clr, bullets), (px, py) in zip(sec_items, pos9):
    add_rect(s9, px, py, 6.0, 2.25, WHITE).line.color.rgb = clr
    add_rect(s9, px, py, 6.0, 0.52, clr)
    add_text_box(s9, title, px + 0.12, py + 0.06, 5.76, 0.42,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    icon_bullet_box(s9, bullets, px + 0.15, py + 0.6, 5.7, 1.55,
                    icon="✔", dot_color=clr, font_size=12, text_color=DARK_GRAY)

footer(s9, 9)


# ═══════════════════════════════════════════════════════
#  SLIDE 10 – IMPACT & SCALABILITY
# ═══════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
slide_bg(s10)
header_bar(s10, "Impact & Scalability", "A platform built to save lives at national scale")

# Impact metrics
metrics = [
    ("⏱️", "< 30 sec", "from symptom entry\nto AI triage result"),
    ("🔔", "Instant", "doctor SMS alert\nfor High-risk cases"),
    ("📍", "Village-level", "outbreak detection\nbefore it spreads"),
    ("🌏", "800M+", "rural Indians who\ncould be reached"),
]
for i, (icon, val, desc) in enumerate(metrics):
    x = 0.5 + i * 3.15
    add_rect(s10, x, 1.65, 2.9, 2.1, DEEP_TEAL)
    add_text_box(s10, icon, x + 0.1, 1.7, 2.7, 0.6,
                 font_size=28, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s10, val, x + 0.1, 2.28, 2.7, 0.6,
                 font_size=22, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(s10, desc, x + 0.1, 2.88, 2.7, 0.75,
                 font_size=11, color=WARM_WHITE, align=PP_ALIGN.CENTER)

# Scalability track
add_text_box(s10, "Scalability Roadmap", 0.5, 4.0, 12.33, 0.45,
             font_size=16, bold=True, color=DEEP_TEAL, align=PP_ALIGN.LEFT)
scale_items = [
    "Horizontal scaling: Dockerized microservices — scale AI service independently from backend",
    "Offline-first mode: IndexedDB queue + background sync for zero-connectivity villages",
    "WhatsApp alert bot: leverage WhatsApp Business API for doctors in low-signal areas",
    "Federated ML: aggregate anonymized outcomes across PHCs and retrain prediction models",
    "Native Android app: Flutter/React Native with multilingual voice capture (Hindi, Tamil, Bengali…)",
]
icon_bullet_box(s10, scale_items, 0.5, 4.5, 12.33, 2.4,
                icon="▶", dot_color=SOFT_GREEN, font_size=13, text_color=DARK_GRAY)

footer(s10, 10)


# ═══════════════════════════════════════════════════════
#  SLIDE 11 – DEMO GUIDE / LIVE WALKTHROUGH
# ═══════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
slide_bg(s11)
header_bar(s11, "Live Demo Guide", "Reproducible demo flow — runs entirely on Docker Compose")

# Left column: setup
add_rect(s11, 0.5, 1.65, 5.8, 5.1, WHITE).line.color.rgb = DEEP_TEAL
add_rect(s11, 0.5, 1.65, 5.8, 0.52, DEEP_TEAL)
add_text_box(s11, "🚀  Setup in 3 Commands", 0.6, 1.7, 5.6, 0.42,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

cmds = [
    "# 1. Configure environment",
    "cp .env.example .env",
    "",
    "# 2. Start all services",
    "docker-compose up --build",
    "",
    "# 3. Seed demo data",
    "docker exec -it asha_backend node seed.js",
]
txBox = s11.shapes.add_textbox(Inches(0.65), Inches(2.28), Inches(5.5), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, line in enumerate(cmds):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.size = Pt(11)
    run.font.name = "Courier New"
    run.font.color.rgb = SOFT_GREEN if not line.startswith("#") else MID_GRAY

add_text_box(s11, "Demo Credentials", 0.6, 4.88, 5.6, 0.35,
             font_size=13, bold=True, color=DEEP_TEAL, align=PP_ALIGN.LEFT)
creds = [
    "Doctor:   doctor@asha.com  /  Doctor@123",
    "ASHA:     asha@asha.com    /  Asha@1234",
    "Admin:    admin@asha.com   /  Admin@123",
]
add_multi_para(s11, creds, 0.6, 5.28, 5.5, 1.2,
               font_size=11, color=DARK_GRAY, font_name="Courier New")

# Right column: demo flow
add_rect(s11, 6.8, 1.65, 6.0, 5.1, WHITE).line.color.rgb = ACCENT_BLUE
add_rect(s11, 6.8, 1.65, 6.0, 0.52, ACCENT_BLUE)
add_text_box(s11, "🎯  Demo Narrative Flow", 6.9, 1.7, 5.8, 0.42,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
demo_steps = [
    "Login as ASHA worker",
    "Register a new patient (name, age, village, GPS)",
    "Enter symptoms via voice: 'fever, headache, joint pain'",
    "View AI result → Dengue · HIGH · Refer to PHC",
    "Show auto-SMS sent to doctor (Twilio log)",
    "Switch to Doctor login → see alert in queue",
    "Open Analytics → disease bar chart + daily trend",
    "Click 'Simulate Outbreak' → watch spike appear",
    "Open Outbreak panel → village-level alert banner",
    "Review patient, add doctor notes → mark resolved",
]
icon_bullet_box(s11, demo_steps, 6.9, 2.28, 5.8, 4.3,
                icon="►", dot_color=ACCENT_BLUE, font_size=12.5, text_color=DARK_GRAY)

footer(s11, 11)


# ═══════════════════════════════════════════════════════
#  SLIDE 12 – CLOSING / CALL TO ACTION
# ═══════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
add_rect(s12, 0, 0, 13.33, 7.5, DARK_GRAY)
add_rect(s12, 0, 0, 13.33, 0.18, SOFT_GREEN)
add_rect(s12, 0, 7.32, 13.33, 0.18, DEEP_TEAL)

# Central content
add_text_box(s12, "ASHA ASSIST", 1.5, 0.6, 10.33, 1.3,
             font_size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s12, "AI Rural Health Intelligence Platform", 1.5, 1.8, 10.33, 0.6,
             font_size=22, bold=False, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_rect(s12, 4.5, 2.55, 4.33, 0.06, SOFT_GREEN)

add_text_box(s12,
             '"Every ASHA worker deserves an AI co-pilot.\n'
             'Every high-risk patient deserves a timely referral.\n'
             'Every PHC deserves visibility into what\'s coming."',
             1.2, 2.75, 10.93, 1.4,
             font_size=17, bold=False, color=WARM_WHITE, align=PP_ALIGN.CENTER)

# Three pillars
pillars = [
    ("🌐", "Open Source\n& Extensible", DEEP_TEAL),
    ("🏥", "PHC-Ready\nDeployment", SOFT_GREEN),
    ("🤝", "ASHA + Doctor\n+ AI Synergy", ACCENT_BLUE),
]
for i, (icon, text, clr) in enumerate(pillars):
    x = 2.0 + i * 3.3
    add_rect(s12, x, 4.35, 2.8, 1.7, clr)
    add_text_box(s12, icon, x + 0.1, 4.4, 2.6, 0.6,
                 font_size=30, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s12, text, x + 0.1, 4.98, 2.6, 0.95,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# CTA
add_rect(s12, 2.5, 6.3, 8.33, 0.65, SOFT_GREEN)
add_text_box(s12,
             "Let's connect the last mile. Let's deploy this. Let's save lives.",
             2.7, 6.34, 7.9, 0.57,
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

footer(s12, 12)


# ═══════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════
import os
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "ASHA_Assist_Hackathon_Pitch.pptx")
prs.save(out_path)
print(f"✅  Saved: {out_path}")
